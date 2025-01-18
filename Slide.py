import os
import re
import pickle
import cv2
import mediapipe as mp
import numpy as np
import speech_recognition as sr
from pptx import Presentation
import keyboard
import time
import pygame
from pynput.keyboard import Key, Controller
from threading import Thread, Event
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import PorterStemmer
import nltk
import sys

nltk.download('stopwords')

# Create an instance of the PorterStemmer
stemmer = PorterStemmer()

# Initialize global variables
current_slide = 1
total_slides = 0
exit_event = Event()  # Event to signal threads to exit

# Function to handle hand swipe detection and key presses
def hand_swipe_detection():
    global current_slide
    global total_slides

    # Load the trained model
    model_dict = pickle.load(open('./modelarr1.p', 'rb')) 
    model = model_dict['model']

    # Initialize the camera
    cap = cv2.VideoCapture(0)

    # Initialize MediaPipe Hands
    mp_hands = mp.solutions.hands
    mp_drawing = mp.solutions.drawing_utils
    mp_drawing_styles = mp.solutions.drawing_styles

    hands = mp_hands.Hands(static_image_mode=False, max_num_hands=2, min_detection_confidence=0.3)

    # Define labels dictionary
    labels_dict = {0: '0', 1: '1', 2: '2', 3: '3', 4: '4', 5: '5'}

    # Initialize keyboard controller
    keyboard_ctrl = Controller()

    # Variable to store previous detected character
    previous_character = None

    while not exit_event.is_set():
        data_aux = []
        x_ = []
        y_ = []

        ret, frame = cap.read()

        if not ret:
            print("Error: Failed to read frame from camera.")
            break

        H, W, _ = frame.shape

        frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)

        results = hands.process(frame_rgb)
        if results.multi_hand_landmarks and results.multi_handedness:
            largest_left_hand_landmarks = None
            max_area = 0

            # Process only left hand
            for idx, handedness in enumerate(results.multi_handedness):
                label = handedness.classification[0].label

                if label == "Left":
                    hand_landmarks = results.multi_hand_landmarks[idx]

                    # Find the largest left hand by calculating the area of its bounding box
                    x_min, y_min = float('inf'), float('inf')
                    x_max, y_max = float('-inf'), float('-inf')

                    for landmark in hand_landmarks.landmark:
                        x = landmark.x
                        y = landmark.y
                        x_min = min(x_min, x)
                        y_min = min(y_min, y)
                        x_max = max(x_max, x)
                        y_max = max(y_max, y)

                    area = (x_max - x_min) * (y_max - y_min)

                    if area > max_area:
                        max_area = area
                        largest_left_hand_landmarks = hand_landmarks

            if largest_left_hand_landmarks:
                # Draw landmarks for the left hand
                mp_drawing.draw_landmarks(
                    frame,  # image to draw on
                    largest_left_hand_landmarks,  # left hand landmarks
                    mp_hands.HAND_CONNECTIONS,  # hand connections
                    mp_drawing_styles.get_default_hand_landmarks_style(),
                    mp_drawing_styles.get_default_hand_connections_style())

                # Extract and normalize coordinates for the left hand
                for i in range(len(largest_left_hand_landmarks.landmark)):
                    x = largest_left_hand_landmarks.landmark[i].x
                    y = largest_left_hand_landmarks.landmark[i].y

                    x_.append(x)
                    y_.append(y)

                for i in range(len(largest_left_hand_landmarks.landmark)):
                    data_aux.append(largest_left_hand_landmarks.landmark[i].x - min(x_))
                    data_aux.append(largest_left_hand_landmarks.landmark[i].y - min(y_))

                # Pad with zeros if necessary to ensure consistent feature length
                if len(data_aux) < 84:
                    data_aux.extend([0] * (84 - len(data_aux)))

                x1 = int(min(x_) * W) - 10
                y1 = int(min(y_) * H) - 10

                x2 = int(max(x_) * W) - 10
                y2 = int(max(y_) * H) - 10

                prediction = model.predict([np.asarray(data_aux)])

                predicted_character = labels_dict[int(prediction[0])]

                cv2.rectangle(frame, (x1, y1), (x2, y2), (0, 0, 0), 4)
                cv2.putText(frame, predicted_character, (x1, y1 - 10), cv2.FONT_HERSHEY_SIMPLEX, 1.3, (0, 0, 0), 3,
                            cv2.LINE_AA)

                # Check for specific sequences and simulate key presses
                if previous_character == '0':
                    if predicted_character == '1':
                        keyboard_ctrl.press(Key.right)
                        keyboard_ctrl.release(Key.right)
                        print("Right arrow key pressed")
                        current_slide = min(current_slide + 1, total_slides)  # Update current slide
                    elif predicted_character == '2':
                        keyboard_ctrl.press(Key.left)
                        keyboard_ctrl.release(Key.left)
                        print("Left arrow key pressed")
                        current_slide = max(current_slide - 1, 1)  # Update current slide

                previous_character = predicted_character

        cv2.imshow('frame', frame)
        if cv2.waitKey(1) & 0xFF == ord('q'):  # Exit on 'q' keystroke
            exit_event.set()
            break

    cap.release()
    cv2.destroyAllWindows()

# Function to handle voice commands
def voice_command_handler():
    def extract_ppt_content(file_path):
        prs = Presentation(file_path)
        slides_content = []

        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + " "
            slides_content.append((i + 1, slide_text.strip()))

        return slides_content, len(prs.slides)

    def find_largest_phrase_slide(ppt_path, search_phrase):
        prs = Presentation(ppt_path)
        slide_composite_sizes = []

        stop_words = set(stopwords.words('english'))
        phrase_words = search_phrase.split()
        filtered_words = [stemmer.stem(word.lower()) for word in phrase_words if word.lower() not in stop_words]

        for slide_index, slide in enumerate(prs.slides):
            max_font_sizes = {word: 0 for word in filtered_words}
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size:
                                    for word in filtered_words:
                                        if stemmer.stem(word.lower()) in run.text.lower():
                                            max_font_sizes[word] = max(max_font_sizes[word], run.font.size.pt)

            composite_size = sum(max_font_sizes.values())
            slide_composite_sizes.append((composite_size, slide_index + 1))

        if slide_composite_sizes:
            largest_composite_slide = max(slide_composite_sizes, key=lambda x: x[0])
            return largest_composite_slide[1]
        else:
            return None

    def listen_to_command():
        recognizer = sr.Recognizer()
        mic = sr.Microphone()

        with mic as source:
            print("Listening for command...")
            recognizer.adjust_for_ambient_noise(source)
            audio = recognizer.listen(source)

        try:
            command = recognizer.recognize_google(audio).lower()
            print(f"You said: {command}")
            return command
        except sr.UnknownValueError:
            print("Sorry, I did not understand that.")
            return None
        except sr.RequestError:
            print("Could not request results from Google Speech Recognition service.")
            return None

    def play_sound(sound_path):
        pygame.mixer.init()
        pygame.mixer.music.load(sound_path)
        pygame.mixer.music.play()
        while pygame.mixer.music.get_busy():
            time.sleep(0.1)

    def navigate_to_slide(current_slide, target_slide, total_slides):
        if target_slide < 1:
            target_slide = 1
        elif target_slide > total_slides:
            target_slide = total_slides

        if target_slide > current_slide:
            for _ in range(target_slide - current_slide):
                keyboard.press_and_release('right')
                time.sleep(0.5)
        elif target_slide < current_slide:
            for _ in range(current_slide - target_slide):
                keyboard.press_and_release('left')
                time.sleep(0.5)

        return target_slide

    global current_slide
    global total_slides

    ppt_path = sys.argv[1]  # Take file path from command line argument
    print(f"Received PPT path in Slide.py: {ppt_path}")
    sound_path = r"C:\Users\vaish\Desktop\Projects\ppt_control_web\venv\bell.wav"

    content, total_slides = extract_ppt_content(ppt_path)
    current_slide = 1

    while not exit_event.is_set():
        command = listen_to_command()
        if command:
            match = re.search(r'hey', command)
            if match:
                play_sound(sound_path)
                topic = command[match.end():].strip()

                if topic:
                    print(f"Searching for slides related to: {topic}")
                    target_slide = find_largest_phrase_slide(ppt_path, topic)
                    if target_slide:
                        print(f"Found relevant slide: {target_slide}")
                        current_slide = navigate_to_slide(current_slide, target_slide, total_slides)
                    else:
                        print("Sorry, no relevant slides found.")
            else:
                print("Command does not contain 'hey', ignoring.")

# Create and start the hand swipe detection and voice command threads
hand_thread = Thread(target=hand_swipe_detection)
voice_thread = Thread(target=voice_command_handler)

hand_thread.start()
voice_thread.start()

# Monitor for 'q' keystroke to end both threads
keyboard.wait('q')
exit_event.set()

# Wait for both threads to finish
hand_thread.join()
voice_thread.join()

print("Exiting program...")
